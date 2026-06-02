"""
Service de matching pitch -> dispositifs de financement.
Flux :
  1. Extraction du texte (PDF / PPTX / TXT)
  2. Nettoyage du bruit documentaire (emails, dates, urls, boilerplate)
  3. Detection d'un profil secteur / pays / type / montants
  4. Matching SQL PostgreSQL avec repli applicatif si besoin
"""

from __future__ import annotations

import io
import json
import re
from datetime import datetime, timezone
from typing import Any

import httpx
from sqlalchemy import BigInteger, Text as SA_Text, bindparam, text
from sqlalchemy.dialects.postgresql import ARRAY as PG_ARRAY
from sqlalchemy.ext.asyncio import AsyncSession
from unidecode import unidecode

from app.config import settings


STOPWORDS = {
    "le", "la", "les", "de", "du", "des", "un", "une", "et", "en", "au", "aux",
    "pour", "par", "sur", "dans", "avec", "est", "sont", "nous", "vous", "ils",
    "qui", "que", "quoi", "dont", "ou", "ce", "cet", "cette", "ces", "mon", "ma",
    "mes", "notre", "votre", "leur", "leurs", "se", "si", "il", "elle", "the",
    "of", "and", "in", "to", "a", "is", "for", "on", "with", "as", "at", "an",
    "be", "by", "this", "that", "are", "from", "it", "we", "our", "its", "mais",
    "donc", "car", "ainsi", "afin", "lors", "plus", "tres", "bien", "tout", "tous",
    "toutes", "aussi", "meme", "comme", "notamment", "contact", "email", "mail",
    "www", "http", "https", "com", "org", "net", "document", "presentation",
    "company", "startup", "october", "january", "february", "march", "april", "may",
    "june", "july", "august", "september", "november", "december", "janvier",
    "fevrier", "mars", "avril", "mai", "juin", "juillet", "aout", "septembre",
    "octobre", "novembre", "decembre",
}

GENERIC_KEYWORDS = {
    "project", "projet", "entreprise", "business", "plan", "pitch", "resume",
    "company", "startup", "document", "presentation", "financial",
}

NOISE_PATTERNS = [
    r"\b[\w\.-]+@[\w\.-]+\.\w+\b",
    r"\bhttps?://\S+\b",
    r"\bwww\.\S+\b",
    r"\b(?:janvier|fevrier|mars|avril|mai|juin|juillet|aout|septembre|octobre|novembre|decembre|"
    r"january|february|march|april|may|june|july|august|september|october|november|december)\b",
    r"\b20\d{2}\b",
    r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b",
]

SECTOR_KEYWORDS = {
    "agriculture": ["agricole", "agriculture", "agri", "elevage", "culture", "recolte", "ferme", "rural"],
    "energie": ["energie", "solaire", "electricite", "renouvelable", "photovoltaique", "biomasse"],
    "sante": ["sante", "medical", "medecine", "hopital", "clinique", "pharmaceutique"],
    "numerique": ["numerique", "digital", "tech", "logiciel", "application", "data", "intelligence", "ia", "saas", "plateforme", "mobile", "payment", "paiement"],
    "education": ["education", "formation", "ecole", "universite", "apprentissage", "pedagogie", "enseignement"],
    "environnement": ["environnement", "ecologie", "climatique", "carbone", "biodiversite", "durable", "vert", "climate"],
    "industrie": ["industrie", "industriel", "manufacture", "production", "usine", "fabrication", "supply", "manufacturing"],
    "transport": ["transport", "logistique", "mobilite", "vehicule", "route", "rail"],
    "finance": ["finance", "fintech", "bancaire", "assurance", "credit", "microfinance", "wallet", "lending", "smb", "sme", "payment", "paiement"],
    "eau": ["eau", "hydraulique", "assainissement", "irrigation", "potable"],
    "social": ["social", "solidaire", "inclusion", "handicap", "precarite", "emploi", "insertion"],
}

COUNTRY_KEYWORDS = {
    "France": ["france", "francais", "hexagone", "metropole"],
    "Sénégal": ["senegal", "senegalais", "dakar"],
    "Côte d'Ivoire": ["cote d'ivoire", "ivoirien", "abidjan"],
    "Maroc": ["maroc", "marocain", "casablanca", "rabat"],
    "Tunisie": ["tunisie", "tunisien", "tunis"],
    "Cameroun": ["cameroun", "camerounais", "yaounde"],
    "Mali": ["mali", "malien", "bamako"],
    "Burkina Faso": ["burkina", "ouagadougou"],
    "Bénin": ["benin", "beninois", "cotonou"],
    "Togo": ["togo", "togolais", "lome"],
    "République démocratique du Congo": ["rdc", "congo", "kinshasa"],
    "Afrique de l'Ouest": ["afrique de l'ouest", "afrique subsaharienne", "cedeao", "ecowas", "uemoa"],
    "Afrique": ["afrique", "africain", "subsaharien", "subsaharienne"],
}

COUNTRY_SCOPE_MAP = {
    "Bénin": ["Afrique", "Afrique de l'Ouest", "International", "France / Afrique"],
    "Burkina Faso": ["Afrique", "Afrique de l'Ouest", "International", "France / Afrique"],
    "Côte d'Ivoire": ["Afrique", "Afrique de l'Ouest", "International", "France / Afrique"],
    "Togo": ["Afrique", "Afrique de l'Ouest", "International", "France / Afrique"],
    "Mali": ["Afrique", "Afrique de l'Ouest", "International", "France / Afrique"],
    "Sénégal": ["Afrique", "Afrique de l'Ouest", "International", "France / Afrique"],
    "Cameroun": ["Afrique", "Afrique centrale", "International", "France / Afrique"],
    "République démocratique du Congo": ["Afrique", "Afrique centrale", "International", "France / Afrique"],
    "Kenya": ["Afrique", "Afrique de l'Est", "International"],
    "Maroc": ["Afrique", "Afrique du Nord", "International"],
    "Tunisie": ["Afrique", "Afrique du Nord", "International"],
    "Afrique de l'Ouest": ["Afrique", "International", "France / Afrique"],
    "Afrique": ["International", "France / Afrique"],
    "France": ["Europe", "International"],
}

DEVICE_TYPE_KEYWORDS = {
    "subvention": ["subvention", "grant", "don", "dotation", "prix", "award", "cash award"],
    "pret": ["pret", "emprunt", "credit", "loan", "remboursable", "debt"],
    "aap": ["appel a projets", "appel a candidatures", "aap", "concours", "soumission"],
    "investissement": ["investissement", "capital", "equity", "fonds propres", "venture", "seed", "series", "investor", "investisseurs", "vc"],
    "accompagnement": ["accompagnement", "mentorat", "incubation", "acceleration", "coaching", "accelerator", "incubator", "incubateur"],
    "garantie": ["garantie", "caution", "aval"],
}

TYPE_INTENT_SIGNALS = {
    "investissement": ["equity", "venture", "vc", "seed", "series", "investor", "investisseurs", "capital", "dilution", "valorisation", "term sheet"],
    "pret": ["loan", "credit", "debt", "emprunt", "pret", "remboursable", "interest", "taux"],
    "subvention": ["grant", "subvention", "don", "dotation", "non dilutif", "non-dilutif", "award"],
    "accompagnement": ["accelerator", "acceleration", "incubation", "coaching", "mentorat", "programme"],
    "garantie": ["garantie", "caution", "aval", "guarantee"],
    "aap": ["appel a projets", "appel a candidatures", "soumission", "concours"],
}

TYPE_CONTRADICTIONS = {
    "investissement": {"subvention", "pret", "garantie", "accompagnement", "aap"},
    "pret": {"investissement", "subvention", "accompagnement", "aap"},
    "subvention": {"investissement", "pret", "garantie"},
    "accompagnement": {"investissement", "pret", "garantie"},
    "garantie": {"investissement", "subvention", "accompagnement", "aap"},
    "aap": {"investissement", "pret", "garantie"},
}

KNOWN_SECTORS = tuple(SECTOR_KEYWORDS.keys())
KNOWN_COUNTRIES = tuple(COUNTRY_KEYWORDS.keys())
KNOWN_DEVICE_TYPES = tuple(DEVICE_TYPE_KEYWORDS.keys())

MATCH_ANALYSIS_SCHEMA: dict[str, Any] = {
    "name": "project_financing_profile",
    "strict": True,
    "schema": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "summary": {"type": "string"},
            "sectors": {"type": "array", "items": {"type": "string", "enum": list(KNOWN_SECTORS)}},
            "countries": {"type": "array", "items": {"type": "string", "enum": list(KNOWN_COUNTRIES)}},
            "types": {"type": "array", "items": {"type": "string", "enum": list(KNOWN_DEVICE_TYPES)}},
            "dominant_type": {"anyOf": [{"type": "string", "enum": list(KNOWN_DEVICE_TYPES)}, {"type": "null"}]},
            "amount_min": {"anyOf": [{"type": "integer"}, {"type": "null"}]},
            "amount_max": {"anyOf": [{"type": "integer"}, {"type": "null"}]},
            "keywords": {"type": "array", "items": {"type": "string"}},
            "project_stage": {"anyOf": [{"type": "string"}, {"type": "null"}]},
            "business_model": {"anyOf": [{"type": "string"}, {"type": "null"}]},
            "target_customers": {"anyOf": [{"type": "string"}, {"type": "null"}]},
            "evidence": {"type": "array", "items": {"type": "string"}},
            "missing_information": {"type": "array", "items": {"type": "string"}},
            "confidence": {"type": "integer", "minimum": 0, "maximum": 100},
        },
        "required": [
            "summary",
            "sectors",
            "countries",
            "types",
            "dominant_type",
            "amount_min",
            "amount_max",
            "keywords",
            "project_stage",
            "business_model",
            "target_customers",
            "evidence",
            "missing_information",
            "confidence",
        ],
    },
}


def _normalize_text(text: str) -> str:
    return unidecode(text.lower())


def _clean_document_text(raw_text: str) -> str:
    cleaned = raw_text
    for pattern in NOISE_PATTERNS:
        cleaned = re.sub(pattern, " ", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"[|•·]", " ", cleaned)
    cleaned = re.sub(r"[ \t]{2,}", " ", cleaned)
    cleaned = re.sub(r"\n{2,}", "\n", cleaned)
    return cleaned.strip()


def _find_matches(text: str, keywords_map: dict[str, list[str]], limit: int) -> list[str]:
    normalized = _normalize_text(text)
    scored: list[tuple[int, str]] = []
    for label, variants in keywords_map.items():
        hits = sum(1 for variant in variants if _normalize_text(variant) in normalized)
        if hits:
            scored.append((hits, label))
    scored.sort(key=lambda item: (-item[0], item[1]))
    return [label for _, label in scored[:limit]]


def _rank_matches(text: str, keywords_map: dict[str, list[str]]) -> list[tuple[str, int]]:
    normalized = _normalize_text(text)
    scored: list[tuple[str, int]] = []
    for label, variants in keywords_map.items():
        hits = sum(1 for variant in variants if _normalize_text(variant) in normalized)
        if hits:
            scored.append((label, hits))
    scored.sort(key=lambda item: (-item[1], item[0]))
    return scored


def _extract_amounts(normalized_text: str) -> tuple[int | None, int | None]:
    amounts: list[int] = []
    for match in re.finditer(
        r"(\d[\d\s]*(?:[.,]\d+)?)\s*(?:millions?|m)?\s*(?:d[' ]*euros?|€|eur|\$|usd|fcfa|xaf|xof)?",
        normalized_text,
    ):
        raw = match.group(1).replace(" ", "").replace(",", ".")
        try:
            value = float(raw)
            if "million" in match.group(0) or re.search(r"\bm\b", match.group(0)):
                value *= 1_000_000
            if 1_000 < value < 100_000_000:
                amounts.append(int(value))
        except ValueError:
            continue
    if not amounts:
        return None, None
    return min(amounts), max(amounts)


def _extract_summary(cleaned_text: str) -> str:
    lines = [line.strip() for line in cleaned_text.splitlines() if len(line.strip()) > 40]
    filtered = [
        line for line in lines
        if "@" not in line and not re.search(r"\b20\d{2}\b", line)
    ]
    source = filtered or lines
    return " ".join(source[:2])[:280]


def analyse_text(raw_text: str) -> dict:
    cleaned_text = _clean_document_text(raw_text)
    normalized_text = _normalize_text(cleaned_text)

    sectors = _find_matches(cleaned_text, SECTOR_KEYWORDS, limit=4)
    countries = _find_matches(cleaned_text, COUNTRY_KEYWORDS, limit=5)
    ranked_types = _rank_matches(cleaned_text, DEVICE_TYPE_KEYWORDS)
    intent_types = _rank_matches(cleaned_text, TYPE_INTENT_SIGNALS)
    combined_type_scores: dict[str, int] = {}
    for label, hits in ranked_types:
        combined_type_scores[label] = combined_type_scores.get(label, 0) + hits * 2
    for label, hits in intent_types:
        combined_type_scores[label] = combined_type_scores.get(label, 0) + hits * 3
    types = [
        label for label, _ in sorted(
            combined_type_scores.items(),
            key=lambda item: (-item[1], item[0]),
        )[:3]
    ]
    dominant_type = types[0] if types else None

    if "fintech" in normalized_text and "finance" not in sectors:
        sectors.insert(0, "finance")
    if any(token in normalized_text for token in ("payment", "paiement", "wallet", "mobile money")) and "numerique" not in sectors:
        sectors.insert(0, "numerique")

    amount_min, amount_max = _extract_amounts(normalized_text)

    words = re.findall(r"[a-z]{4,}", normalized_text)
    keywords = [
        word for word in words
        if word not in STOPWORDS
        and word not in GENERIC_KEYWORDS
        and not word.isdigit()
    ]

    freq: dict[str, int] = {}
    for word in keywords:
        freq[word] = freq.get(word, 0) + 1

    top_keywords = sorted(freq, key=lambda item: (-freq[item], -len(item), item))[:20]
    summary = _extract_summary(cleaned_text)

    return {
        "sectors": sectors[:4],
        "countries": countries[:5],
        "country_scopes": _country_scope_values(countries[:5]),
        "regional_scopes": _regional_scope_values(countries[:5]),
        "geographic_scopes": _geographic_scope_values(countries[:5]),
        "types": types[:3],
        "dominant_type": dominant_type,
        "amount_min": amount_min,
        "amount_max": amount_max,
        "keywords": top_keywords,
        "summary": summary,
        "query": " ".join(top_keywords[:12]),
    }


def _clip_document_for_ai(raw_text: str, max_chars: int = 18000) -> str:
    cleaned = _clean_document_text(raw_text)
    if len(cleaned) <= max_chars:
        return cleaned
    head = cleaned[: int(max_chars * 0.7)]
    tail = cleaned[-int(max_chars * 0.3) :]
    return f"{head}\n\n[...document tronque pour analyse...]\n\n{tail}"


def _country_scope_values(countries: list[str]) -> list[str]:
    values: list[str] = []
    for country in countries or []:
        for value in [country, *COUNTRY_SCOPE_MAP.get(country, [])]:
            if value and value not in values:
                values.append(value)
    return values


def _regional_scope_values(countries: list[str]) -> list[str]:
    return [
        value
        for value in _country_scope_values(countries)
        if value not in {"International", "France / Afrique"}
    ]


def _geographic_scope_values(countries: list[str]) -> list[str]:
    if not countries:
        return []
    return ["continental", "multi_pays", "international"]


def _dedupe_known(values: list[Any], allowed: tuple[str, ...], limit: int) -> list[str]:
    allowed_by_norm = {_normalize_text(value): value for value in allowed}
    result: list[str] = []
    for raw in values or []:
        value = allowed_by_norm.get(_normalize_text(str(raw)))
        if value and value not in result:
            result.append(value)
    return result[:limit]


def _normalize_chatgpt_profile(payload: dict[str, Any], heuristic: dict) -> dict:
    sectors = _dedupe_known(payload.get("sectors") or [], KNOWN_SECTORS, 4)
    countries = _dedupe_known(payload.get("countries") or [], KNOWN_COUNTRIES, 5)
    types = _dedupe_known(payload.get("types") or [], KNOWN_DEVICE_TYPES, 3)

    if not sectors:
        sectors = heuristic.get("sectors", [])[:4]
    if not countries:
        countries = heuristic.get("countries", [])[:5]
    if not types:
        types = heuristic.get("types", [])[:3]

    dominant_type = payload.get("dominant_type")
    if dominant_type not in KNOWN_DEVICE_TYPES:
        dominant_type = types[0] if types else heuristic.get("dominant_type")
    if dominant_type and dominant_type not in types:
        types = [dominant_type, *[item for item in types if item != dominant_type]][:3]

    amount_min = payload.get("amount_min")
    amount_max = payload.get("amount_max")
    if not isinstance(amount_min, int):
        amount_min = heuristic.get("amount_min")
    if not isinstance(amount_max, int):
        amount_max = heuristic.get("amount_max")
    if isinstance(amount_min, int) and isinstance(amount_max, int) and amount_min > amount_max:
        amount_min, amount_max = amount_max, amount_min

    raw_keywords = [str(item).strip().lower() for item in payload.get("keywords") or []]
    keywords = [
        keyword
        for keyword in raw_keywords
        if keyword and keyword not in STOPWORDS and keyword not in GENERIC_KEYWORDS
    ][:20]
    if len(keywords) < 5:
        keywords = list(dict.fromkeys([*keywords, *heuristic.get("keywords", [])]))[:20]

    confidence = payload.get("confidence")
    if not isinstance(confidence, int):
        confidence = 70
    confidence = max(0, min(100, confidence))

    summary = str(payload.get("summary") or heuristic.get("summary") or "").strip()
    if len(summary) > 420:
        summary = summary[:417].rstrip() + "..."

    profile = {
        "sectors": sectors,
        "countries": countries,
        "country_scopes": _country_scope_values(countries),
        "regional_scopes": _regional_scope_values(countries),
        "geographic_scopes": _geographic_scope_values(countries),
        "types": types,
        "dominant_type": dominant_type,
        "amount_min": amount_min,
        "amount_max": amount_max,
        "keywords": keywords,
        "summary": summary,
        "query": " ".join(keywords[:12]) or heuristic.get("query") or "financement projet",
        "project_stage": payload.get("project_stage"),
        "business_model": payload.get("business_model"),
        "target_customers": payload.get("target_customers"),
        "evidence": [str(item).strip()[:220] for item in payload.get("evidence") or [] if str(item).strip()][:6],
        "missing_information": [
            str(item).strip()[:160] for item in payload.get("missing_information") or [] if str(item).strip()
        ][:8],
        "confidence": confidence,
        "analysis_provider": "openai",
        "analysis_model": settings.AI_MATCH_MODEL,
    }
    return profile


def build_match_analysis_prompt(raw_text: str, heuristic: dict) -> str:
    allowed = {
        "sectors": list(KNOWN_SECTORS),
        "countries": list(KNOWN_COUNTRIES),
        "types": list(KNOWN_DEVICE_TYPES),
    }
    return (
        "Analyse ce document de projet pour une plateforme de matching financement.\n"
        "Objectif: extraire un profil de financement exact, exploitable et prudent.\n\n"
        "Regles de precision:\n"
        "- Base-toi uniquement sur le document fourni, sans inventer de pays, secteur, montant ou stade.\n"
        "- Si une information n'est pas clairement presente, laisse null ou ajoute-la a missing_information.\n"
        "- Choisis uniquement des valeurs parmi les listes autorisees.\n"
        "- Le dominant_type doit exprimer le besoin principal de financement, pas tous les mots cites.\n"
        "- Pour les montants, convertis en entiers approximatifs dans la devise mentionnee; si incertain, null.\n"
        "- Les evidence doivent etre des indices courts tires ou paraphrases du document.\n"
        "- Les keywords doivent etre specifiques au projet, pas generiques.\n\n"
        f"Listes autorisees:\n{json.dumps(allowed, ensure_ascii=False)}\n\n"
        f"Premiere detection locale a verifier/corriger:\n{json.dumps(heuristic, ensure_ascii=False, default=str)}\n\n"
        f"Document:\n{_clip_document_for_ai(raw_text)}"
    )


async def _call_openai_match_analysis(raw_text: str, heuristic: dict) -> dict[str, Any]:
    if not settings.OPENAI_API_KEY:
        raise RuntimeError("OPENAI_API_KEY manquante")

    headers = {
        "Authorization": f"Bearer {settings.OPENAI_API_KEY}",
        "Content-Type": "application/json",
    }
    body = {
        "model": settings.AI_MATCH_MODEL,
        "messages": [
            {
                "role": "system",
                "content": (
                    "Tu es un analyste financement senior. Tu extrais des profils projet fiables "
                    "pour matcher des aides, subventions, prets, garanties, concours et investisseurs. "
                    "Tu es strict: aucune supposition non etayee."
                ),
            },
            {"role": "user", "content": build_match_analysis_prompt(raw_text, heuristic)},
        ],
        "temperature": 0,
        "response_format": {"type": "json_schema", "json_schema": MATCH_ANALYSIS_SCHEMA},
    }
    async with httpx.AsyncClient(timeout=settings.AI_MATCH_TIMEOUT_SECONDS) as client:
        response = await client.post("https://api.openai.com/v1/chat/completions", headers=headers, json=body)
        response.raise_for_status()
        data = response.json()

    content = data["choices"][0]["message"]["content"]
    return json.loads(content)


async def analyse_text_with_chatgpt(raw_text: str) -> dict:
    heuristic = analyse_text(raw_text)
    payload = await _call_openai_match_analysis(raw_text, heuristic)
    return _normalize_chatgpt_profile(payload, heuristic)


MATCH_SQL = """
WITH scored AS (
    SELECT
        d.id,
        d.title,
        d.short_description AS description_courte,
        d.device_type,
        d.country,
        d.region,
        d.zone,
        d.geographic_scope,
        d.sectors,
        d.amount_min,
        d.amount_max,
        d.source_url,
        d.close_date,
        d.validation_status,
        s.reliability AS source_reliability,
        s.consecutive_errors AS source_errors,
        s.is_active AS source_is_active,
        s.last_success_at AS source_last_success_at,
        d.last_verified_at,
        CASE
            WHEN d.search_vector @@ plainto_tsquery('french', :query)
                THEN ts_rank(d.search_vector, plainto_tsquery('french', :query)) * 100
            ELSE 0
        END
        + CASE WHEN d.sectors && :sectors THEN 20 ELSE 0 END
        + CASE WHEN d.country = ANY(:countries) THEN 30 ELSE 0 END
        + CASE
            WHEN d.country = ANY(:country_scopes)
              OR COALESCE(d.region, '') = ANY(:regional_scopes)
              OR COALESCE(d.zone, '') = ANY(:regional_scopes)
              OR COALESCE(d.geographic_scope, '') = ANY(:geographic_scopes)
                THEN 18
            ELSE 0
        END
        + CASE WHEN d.device_type = ANY(:types) THEN 12 ELSE 0 END
        + CASE WHEN :dominant_type IS NOT NULL AND d.device_type = :dominant_type THEN 18 ELSE 0 END
        + CASE
            WHEN :dominant_type = 'investissement' AND d.device_type IN ('subvention', 'pret', 'garantie', 'accompagnement', 'aap') THEN -18
            WHEN :dominant_type = 'pret' AND d.device_type IN ('investissement', 'subvention', 'accompagnement', 'aap') THEN -18
            WHEN :dominant_type = 'subvention' AND d.device_type IN ('investissement', 'pret', 'garantie') THEN -18
            WHEN :dominant_type = 'accompagnement' AND d.device_type IN ('investissement', 'pret', 'garantie') THEN -18
            WHEN :dominant_type = 'garantie' AND d.device_type IN ('investissement', 'subvention', 'accompagnement', 'aap') THEN -18
            WHEN :dominant_type = 'aap' AND d.device_type IN ('investissement', 'pret', 'garantie') THEN -18
            ELSE 0
        END
        + CASE
            WHEN (CAST(:amount_max AS BIGINT) IS NULL OR d.amount_min IS NULL OR d.amount_min <= CAST(:amount_max AS BIGINT))
             AND (CAST(:amount_min AS BIGINT) IS NULL OR d.amount_max IS NULL OR d.amount_max >= CAST(:amount_min AS BIGINT))
                THEN 5
            ELSE 0
        END
        + CASE
            WHEN COALESCE(s.reliability, 3) >= 5 THEN 8
            WHEN COALESCE(s.reliability, 3) = 4 THEN 4
            WHEN COALESCE(s.reliability, 3) <= 2 THEN -6
            ELSE 0
        END
        + CASE
            WHEN COALESCE(s.consecutive_errors, 0) >= 5 THEN -18
            WHEN COALESCE(s.consecutive_errors, 0) >= 3 THEN -12
            WHEN COALESCE(s.consecutive_errors, 0) >= 1 THEN -5
            ELSE 0
        END
        + CASE
            WHEN s.is_active = false THEN -10
            ELSE 0
        END
        + CASE
            WHEN s.last_success_at IS NULL AND d.last_verified_at IS NULL THEN -12
            WHEN s.last_success_at IS NOT NULL AND s.last_success_at >= NOW() - INTERVAL '45 days' THEN 6
            WHEN s.last_success_at IS NOT NULL AND s.last_success_at < NOW() - INTERVAL '120 days' THEN -8
            WHEN d.last_verified_at IS NOT NULL AND d.last_verified_at >= NOW() - INTERVAL '30 days' THEN 4
            WHEN d.last_verified_at IS NOT NULL AND d.last_verified_at < NOW() - INTERVAL '120 days' THEN -6
            ELSE 0
        END AS score
    FROM devices d
    LEFT JOIN sources s ON s.id = d.source_id
    WHERE d.validation_status IN ('approved', 'validated', 'auto_published')
      AND (
          d.search_vector @@ plainto_tsquery('french', :query)
          OR d.sectors && :sectors
          OR d.country = ANY(:countries)
          OR d.country = ANY(:country_scopes)
          OR COALESCE(d.region, '') = ANY(:regional_scopes)
          OR COALESCE(d.zone, '') = ANY(:regional_scopes)
          OR COALESCE(d.geographic_scope, '') = ANY(:geographic_scopes)
          OR d.device_type = ANY(:types)
      )
      AND (
          CARDINALITY(:countries) = 0
          OR d.country = ANY(:countries)
          OR d.country = ANY(:country_scopes)
          OR COALESCE(d.region, '') = ANY(:regional_scopes)
          OR COALESCE(d.zone, '') = ANY(:regional_scopes)
          OR COALESCE(d.geographic_scope, '') = ANY(:geographic_scopes)
      )
)
SELECT * FROM scored
WHERE score > 0
ORDER BY score DESC
LIMIT :limit
"""

_MATCH_STMT = text(MATCH_SQL).bindparams(
    bindparam("sectors", type_=PG_ARRAY(SA_Text)),
    bindparam("countries", type_=PG_ARRAY(SA_Text)),
    bindparam("country_scopes", type_=PG_ARRAY(SA_Text)),
    bindparam("regional_scopes", type_=PG_ARRAY(SA_Text)),
    bindparam("geographic_scopes", type_=PG_ARRAY(SA_Text)),
    bindparam("types", type_=PG_ARRAY(SA_Text)),
    bindparam("dominant_type", type_=SA_Text),
    bindparam("amount_min", type_=BigInteger),
    bindparam("amount_max", type_=BigInteger),
)

FALLBACK_MATCH_SQL = """
SELECT
    d.id,
    d.title,
    d.short_description AS description_courte,
    d.device_type,
    d.country,
    d.region,
    d.zone,
    d.geographic_scope,
    d.sectors,
    d.amount_min,
    d.amount_max,
    d.source_url,
    d.close_date,
    d.validation_status,
    d.last_verified_at,
    s.reliability AS source_reliability,
    s.consecutive_errors AS source_errors,
    s.is_active AS source_is_active,
    s.last_success_at AS source_last_success_at
FROM devices d
LEFT JOIN sources s ON s.id = d.source_id
WHERE d.validation_status IN ('approved', 'validated', 'auto_published')
ORDER BY d.updated_at DESC
LIMIT 300
"""

_FALLBACK_MATCH_STMT = text(FALLBACK_MATCH_SQL)


def _score_fallback_row(row: dict, profile: dict) -> int:
    haystack = " ".join(
        [
            str(row.get("title") or ""),
            str(row.get("description_courte") or ""),
            " ".join(row.get("sectors") or []),
            str(row.get("country") or ""),
            str(row.get("device_type") or ""),
        ]
    ).lower()

    score = 0
    for keyword in profile.get("keywords", [])[:12]:
        if keyword and keyword.lower() in haystack:
            score += 8

    country_scopes = set(profile.get("country_scopes") or _country_scope_values(profile.get("countries") or []))
    regional_scopes = set(profile.get("regional_scopes") or _regional_scope_values(profile.get("countries") or []))
    geographic_scopes = set(profile.get("geographic_scopes") or _geographic_scope_values(profile.get("countries") or []))
    has_country_filter = bool(profile.get("countries"))
    has_geo_match = (
        row.get("country") in (profile.get("countries") or [])
        or row.get("country") in country_scopes
        or row.get("region") in regional_scopes
        or row.get("zone") in regional_scopes
        or row.get("geographic_scope") in geographic_scopes
    )
    if has_country_filter and not has_geo_match:
        return -999

    if row.get("country") in (profile.get("countries") or []):
        score += 30
    elif has_geo_match:
        score += 18
    if row.get("device_type") in (profile.get("types") or []):
        score += 12
    if profile.get("dominant_type") and row.get("device_type") == profile.get("dominant_type"):
        score += 18
    if (
        profile.get("dominant_type")
        and row.get("device_type") in TYPE_CONTRADICTIONS.get(profile.get("dominant_type"), set())
    ):
        score -= 18

    row_sectors = set(row.get("sectors") or [])
    profile_sectors = set(profile.get("sectors") or [])
    score += min(20, 10 * len(row_sectors & profile_sectors))

    amount_min = profile.get("amount_min")
    amount_max = profile.get("amount_max")
    row_amount_min = row.get("amount_min")
    row_amount_max = row.get("amount_max")
    if (
        (amount_max is None or row_amount_min is None or row_amount_min <= amount_max)
        and (amount_min is None or row_amount_max is None or row_amount_max >= amount_min)
    ):
        score += 5

    source_reliability = row.get("source_reliability") or 3
    if source_reliability >= 5:
        score += 8
    elif source_reliability == 4:
        score += 4
    elif source_reliability <= 2:
        score -= 6

    source_errors = row.get("source_errors") or 0
    if source_errors >= 5:
        score -= 18
    elif source_errors >= 3:
        score -= 12
    elif source_errors >= 1:
        score -= 5

    if row.get("source_is_active") is False:
        score -= 10

    source_last_success = row.get("source_last_success_at")
    device_last_verified = row.get("last_verified_at")
    if source_last_success is None and device_last_verified is None:
        score -= 12
    else:
        now = datetime.now(timezone.utc)
        freshest = None
        for candidate in (source_last_success, device_last_verified):
            if candidate is None:
                continue
            if getattr(candidate, "tzinfo", None) is None:
                candidate = candidate.replace(tzinfo=timezone.utc)
            if freshest is None or candidate > freshest:
                freshest = candidate
        if freshest is not None:
            age_days = max((now - freshest).days, 0)
            if age_days <= 45:
                score += 6
            elif age_days >= 120:
                score -= 8

    return score


async def _find_matching_devices_fallback(db: AsyncSession, profile: dict, limit: int) -> list[dict]:
    result = await db.execute(_FALLBACK_MATCH_STMT)
    rows = result.mappings().all()

    devices = []
    for row in rows:
        device = dict(row)
        raw_score = _score_fallback_row(device, profile)
        if raw_score <= 0:
            continue
        device["match_score"] = min(99, raw_score)
        devices.append(device)

    devices.sort(key=lambda item: item["match_score"], reverse=True)
    return devices[:limit]


async def find_matching_devices(db: AsyncSession, profile: dict, limit: int = 15) -> list[dict]:
    query_text = profile.get("query") or " ".join(profile.get("keywords", []))
    sectors = profile.get("sectors") or []
    countries = profile.get("countries") or []
    country_scopes = profile.get("country_scopes") or _country_scope_values(countries)
    regional_scopes = profile.get("regional_scopes") or _regional_scope_values(countries)
    geographic_scopes = profile.get("geographic_scopes") or _geographic_scope_values(countries)
    types = profile.get("types") or []
    dominant_type = profile.get("dominant_type")
    amount_min = profile.get("amount_min")
    amount_max = profile.get("amount_max")

    if not query_text and not sectors and not countries:
        return []

    try:
        result = await db.execute(
            _MATCH_STMT,
            {
                "query": query_text or "financement projet",
                "sectors": sectors,
                "countries": countries,
                "country_scopes": country_scopes,
                "regional_scopes": regional_scopes,
                "geographic_scopes": geographic_scopes,
                "types": types,
                "dominant_type": dominant_type,
                "amount_min": amount_min,
                "amount_max": amount_max,
                "limit": limit,
            },
        )
        rows = result.mappings().all()

        devices = []
        for row in rows:
            device = dict(row)
            raw_score = float(device.pop("score", 0))
            device["match_score"] = min(99, round(raw_score / 1.5))
            devices.append(device)
        return devices
    except Exception:
        return await _find_matching_devices_fallback(db, profile, limit)


def extract_text_from_pdf(content: bytes) -> str:
    try:
        import pdfplumber

        parts: list[str] = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                text_content = page.extract_text()
                if text_content:
                    parts.append(text_content)
        result = "\n".join(parts)
        if not result.strip():
            raise ValueError("Le PDF ne contient pas de texte extractible.")
        return result
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError(f"Impossible de lire le PDF : {exc}") from exc


def extract_text_from_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(content))
        parts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        result = "\n".join(parts)
        if not result.strip():
            raise ValueError("La presentation ne contient pas de texte extractible.")
        return result
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError(f"Impossible de lire le fichier PPTX : {exc}") from exc


def extract_text(filename: str, content: bytes) -> str:
    extension = filename.rsplit(".", 1)[-1].lower()
    if extension == "pdf":
        return extract_text_from_pdf(content)
    if extension in ("pptx", "ppt"):
        return extract_text_from_pptx(content)
    return content.decode("utf-8", errors="ignore")


async def match_project(
    db: AsyncSession,
    filename: str,
    content: bytes,
    limit: int = 15,
) -> dict:
    raw_text = extract_text(filename, content)
    if len(raw_text.strip()) < 30:
        raise ValueError("Le document semble vide ou illisible.")

    try:
        profile = await analyse_text_with_chatgpt(raw_text)
    except Exception as exc:
        raise RuntimeError(f"Analyse ChatGPT indisponible: {type(exc).__name__}") from exc
    matches = await find_matching_devices(db, profile, limit=limit)

    return {
        "profile": {
            "sectors": profile["sectors"],
            "countries": profile["countries"],
            "country_scopes": profile.get("country_scopes", []),
            "types": profile["types"],
            "dominant_type": profile["dominant_type"],
            "amount_min": profile["amount_min"],
            "amount_max": profile["amount_max"],
            "keywords": profile["keywords"][:8],
            "summary": profile["summary"],
            "project_stage": profile.get("project_stage"),
            "business_model": profile.get("business_model"),
            "target_customers": profile.get("target_customers"),
            "evidence": profile.get("evidence", []),
            "missing_information": profile.get("missing_information", []),
            "confidence": profile.get("confidence"),
            "analysis_provider": profile.get("analysis_provider"),
            "analysis_model": profile.get("analysis_model"),
        },
        "matches": matches,
        "total": len(matches),
    }
